import fitz  # PyMuPDF for PDFs
import docx
import requests
from bs4 import BeautifulSoup
from tempfile import NamedTemporaryFile
import math
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor

#for various File Types:
import openpyxl
from PIL import Image
import pytesseract
from pptx import Presentation
import zipfile
import io
from urllib.parse import urlparse
import mimetypes
import os

# def extract_text_from_url(file_url: str) -> str:
#     response = requests.get(file_url)
#     ext = file_url.split('?')[0].split('.')[-1].lower()

#     with NamedTemporaryFile(delete=False, suffix=f".{ext}") as f:
#         f.write(response.content)
#         f.flush()

#         if ext == "pdf":
#             doc = fitz.open(f.name)
#             return "\n".join([page.get_text() for page in doc])

#         elif ext == "docx":
#             doc = docx.Document(f.name)
#             return "\n".join([p.text for p in doc.paragraphs])

#         elif ext == "eml":
#             with open(f.name, "r", encoding="utf-8", errors="ignore") as email_file:
#                 html = email_file.read()
#                 soup = BeautifulSoup(html, "html.parser")
#                 return soup.get_text(separator="\n")

#         else:
#             return "❌ Unsupported file format"

def extract_text_from_url(file_url: str) -> dict:
    """Extract text from various file formats with security checks"""
    
    # Check for fraudulent URLs
    if is_fraudulent_url(file_url):
        return {
            "isError": True,
            "message": "Suspicious or fraudulent URL detected. File processing blocked for security reasons."
        }
    
    try:
        # Get file info without downloading
        head_response = requests.head(file_url, timeout=10)
        file_size = get_file_size(head_response)
        
        # Check file size (1GB limit)
        if file_size > 1024 * 1024 * 1024:  # 1GB in bytes
            return {
                "isError": True,
                "message": f"File too large ({file_size / (1024*1024*1024):.2f} GB). Maximum allowed size is 1GB."
            }
        
        # Download with streaming and size check
        response = requests.get(file_url, stream=True, timeout=30)
        response.raise_for_status()
        
        # Double-check file size during download
        downloaded_size = 0
        content_chunks = []
        
        for chunk in response.iter_content(chunk_size=8192):
            downloaded_size += len(chunk)
            if downloaded_size > 1024 * 1024 * 1024:  # 1GB limit
                return {
                    "isError": True,
                    "message": "File size exceeded during download. Processing stopped for security reasons."
                }
            content_chunks.append(chunk)
        
        file_content = b''.join(content_chunks)
        
        # Determine file extension
        ext = file_url.split('?')[0].split('.')[-1].lower() if '.' in file_url else ''
        
        # If no extension in URL, try to detect from content-type
        if not ext:
            content_type = response.headers.get('content-type', '').lower()
            ext_map = {
                'application/pdf': 'pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                'image/jpeg': 'jpg',
                'image/png': 'png',
                'application/zip': 'zip',
                'text/html': 'eml'
            }
            ext = ext_map.get(content_type, '')
        
        if not ext:
            return {
                "isError": True,
                "message": "Unable to determine file type. Please ensure the file has a valid extension."
            }
        
        # Create temporary file
        with NamedTemporaryFile(delete=False, suffix=f".{ext}") as f:
            f.write(file_content)
            f.flush()
            
            try:
                extracted_text = ""
                
                if ext == "pdf":
                    doc = fitz.open(f.name)
                    extracted_text = "\n".join([page.get_text() for page in doc])

                elif ext == "docx":
                    doc = docx.Document(f.name)
                    extracted_text = "\n".join([p.text for p in doc.paragraphs])

                elif ext == "xlsx":
                    excel_result = extract_text_from_xlsx(f.name)
                    print(f"{excel_result} From extract_text_fromURL")
                    if excel_result.get("isError"):
                        return excel_result
                    
                    extracted_text = excel_result["text"]
                    # Store excel_data for later use in chunking
                    return {
                        "isError": False,
                        "text": extracted_text,
                        "excel_data": excel_result.get("excel_data"),
                        "file_type": "xlsx"
                    }

                elif ext in ["jpg", "jpeg", "png"]:
                    extracted_text = extract_text_from_image(f.name)

                elif ext == "pptx":
                    extracted_text = extract_text_from_pptx(f.name)

                elif ext == "zip":
                    extracted_text = extract_text_from_zip(f.name)

                elif ext == "eml":
                    with open(f.name, "r", encoding="utf-8", errors="ignore") as email_file:
                        html = email_file.read()
                        soup = BeautifulSoup(html, "html.parser")
                        extracted_text = soup.get_text(separator="\n")

                else:
                    return {
                        "isError": True,
                        "message": f"Unsupported file format: .{ext}. Supported formats: PDF, DOCX, XLSX, JPG, PNG, PPTX, ZIP, EML"
                    }
                
                # Check if extraction was successful
                if not extracted_text or extracted_text.strip() == "":
                    return {
                        "isError": True,
                        "message": f"No content could be extracted from the {ext.upper()} file. The file might be empty or corrupted."
                    }
                
                return {
                    "isError": False,
                    "text": extracted_text
                }
                    
            finally:
                # Clean up temporary file
                try:
                    os.unlink(f.name)
                except:
                    pass
    
    except requests.exceptions.Timeout:
        return {
            "isError": True,
            "message": "Request timeout. The file took too long to download. Please try again or use a different file."
        }
    except requests.exceptions.RequestException as e:
        return {
            "isError": True,
            "message": f"Network error occurred while downloading the file: {str(e)}"
        }
    except Exception as e:
        return {
            "isError": True,
            "message": f"Unexpected error occurred while processing the file: {str(e)}"
        }

# Add this to your processor.py file

def chunk_text_enhanced(text: str, chunk_size: int = 300, overlap: int = 50) -> list:
    """Enhanced chunking that matches your working standalone code"""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

def chunk_text_parallel_enhanced(text: str, chunk_size: int = 300, overlap: int = 50, num_threads: int = 4) -> list:
    """
    Enhanced parallel chunking with better parameters
    """
    words = text.split()
    total_words = len(words)
    
    # If document is small, use enhanced regular chunking
    if total_words < 1000:
        return chunk_text_enhanced(text, chunk_size, overlap)
    
    # Calculate words per thread
    words_per_thread = math.ceil(total_words / num_threads)
    
    def process_segment(start_idx, end_idx):
        segment_words = words[start_idx:end_idx]
        segment_text = " ".join(segment_words)
        return chunk_text_enhanced(segment_text, chunk_size, overlap)
    
    # Create thread segments with overlap to maintain context
    segments = []
    for i in range(num_threads):
        start_idx = i * words_per_thread
        end_idx = min((i + 1) * words_per_thread + overlap, total_words)
        if start_idx < total_words:
            segments.append((start_idx, end_idx))
    
    # Process segments in parallel
    all_chunks = []
    with ThreadPoolExecutor(max_workers=num_threads) as executor:
        futures = [executor.submit(process_segment, start, end) for start, end in segments]
        
        for future in concurrent.futures.as_completed(futures):
            chunks = future.result()
            all_chunks.extend(chunks)
    
    return all_chunks


def is_fraudulent_url(url: str) -> bool:
    """Check if URL appears to be fraudulent or suspicious"""
    suspicious_patterns = [
        '.bin', '.exe', '.bat', '.cmd', '.scr', '.com', '.pif',
        'download', 'temp', 'trash', 'dummy', 'test-file',
        'speed-test', 'bandwidth', 'hetzner.com', 'speedtest'
    ]
    
    url_lower = url.lower()
    return any(pattern in url_lower for pattern in suspicious_patterns)

def get_file_size(response) -> int:
    """Get file size from response headers"""
    content_length = response.headers.get('content-length')
    if content_length:
        return int(content_length)
    return 0

def extract_text_from_xlsx(file_path: str) -> dict:
    """Extract text from Excel file separating structured data from supporting text"""
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        excel_data = {
            "structured_rows": [],
            "supporting_text": [],
            "is_excel": True
        }
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            
            # Find structured data table
            table_start_row, headers = find_table_start(sheet)
            
            if table_start_row is None:
                # No structured table found, treat all as supporting text
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_text):
                        excel_data["supporting_text"].append(" | ".join(row_text))
                continue
            
            # Extract supporting text (everything before structured table)
            for row_idx in range(1, table_start_row):
                row_values = []
                for cell in sheet[row_idx]:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                
                if row_values and any(val.strip() for val in row_values):
                    excel_data["supporting_text"].append(" | ".join(row_values))
            
            # Extract structured data (from table_start_row onwards)
            for row_idx, row in enumerate(sheet.iter_rows(min_row=table_start_row + 1, values_only=True)):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                
                if any(row_data) and not is_empty_or_header_row(row_data, headers):
                    # Create row text with headers for context - FIX HERE
                    row_with_headers = []
                    for header, value in zip(headers, row_data):
                        if value.strip():
                            row_with_headers.append(f"{header}: {value}")
                    
                    if row_with_headers:
                        # ENSURE row_text is a STRING, not a list
                        row_text_string = " | ".join(row_with_headers)  # Convert to string
                        
                        excel_data["structured_rows"].append({
                            "sheet": sheet_name,
                            "row_index": table_start_row + row_idx + 1,
                            "row_text": row_text_string,  # This should be a string
                            "raw_values": row_data,
                            "headers": headers  # Add headers for easier access
                        })
        
        # Create main text from structured data only
        main_text_parts = []
        current_sheet = None
        for row in excel_data["structured_rows"]:
            if current_sheet != row["sheet"]:
                current_sheet = row["sheet"]
                main_text_parts.append(f"\nSheet: {current_sheet}")
            main_text_parts.append(row["row_text"])
        
        main_text = "\n".join(main_text_parts) if main_text_parts else ""
        
        # Combine supporting text
        supporting_text_list = excel_data.get("supporting_text", [])
        supporting_text = "\n".join(supporting_text_list) if supporting_text_list else ""
        
        return {
            "isError": False,
            "text": main_text,
            "supporting_text": supporting_text,
            "excel_data": excel_data
        }
        
    except Exception as e:
        return {
            "isError": True,
            "message": f"Failed to read Excel file: {str(e)}"
        }

def extract_text_from_image(file_path: str) -> str:
    """Extract text from image using OCR"""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text.strip() if text.strip() else "No text found in the image"
    except Exception as e:
        raise Exception(f"Failed to process image file: {str(e)}")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file including text from images using OCR"""
    try:
        presentation = Presentation(file_path)
        text_content = []
        
        for i, slide in enumerate(presentation.slides, 1):
            text_content.append(f"Slide {i}:")
            slide_text = []
            slide_image_text = []
            
            for shape in slide.shapes:
                # Extract regular text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from images using OCR
                elif shape.shape_type == 13:  # Picture shape type
                    temp_img_path = None
                    try:
                        # Get image data from shape
                        image_data = shape.image.blob
                        
                        # Create temporary image file
                        with NamedTemporaryFile(delete=False, suffix=".png") as temp_img:
                            temp_img.write(image_data)
                            temp_img.flush()
                            temp_img_path = temp_img.name
                        
                        # Extract text using OCR (file is now closed)
                        image = Image.open(temp_img_path)
                        ocr_text = pytesseract.image_to_string(image, config='--psm 6')
                        image.close()  # Close the image explicitly
                        
                        if ocr_text.strip():
                            slide_image_text.append(f"[Image Text]: {ocr_text.strip()}")
                        
                    except Exception as ocr_error:
                        print(f"⚠️ OCR failed for image in slide {i}: {ocr_error}")
                    
                    finally:
                        # Clean up temp file safely
                        if temp_img_path and os.path.exists(temp_img_path):
                            try:
                                os.unlink(temp_img_path)
                            except PermissionError:
                                # If file is still locked, try to delete later
                                import atexit
                                atexit.register(lambda: os.unlink(temp_img_path) if os.path.exists(temp_img_path) else None)
            
            # Add all extracted text for this slide
            if slide_text:
                text_content.extend(slide_text)
            if slide_image_text:
                text_content.extend(slide_image_text)
            
            text_content.append("")  # Empty line between slides
        
        result = "\n".join(text_content)
        slide_count = len(presentation.slides)
        print(f"📊 DEBUG: Extracted text from {slide_count} slides from PowerPoint (including OCR)")
        return result if result.strip() else "PowerPoint file appears to contain no text"
        
    except Exception as e:
        raise Exception(f"Failed to read PowerPoint file: {str(e)}")

def extract_text_from_zip(file_path: str) -> str:
    """Extract text from ZIP file (list contents and extract text files)"""
    try:
        text_content = ["ZIP Archive Contents:"]
        
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            file_list = zip_file.namelist()
            text_content.extend([f"- {filename}" for filename in file_list[:50]])  # Limit to 50 files
            
            if len(file_list) > 50:
                text_content.append(f"... and {len(file_list) - 50} more files")
            
            # Try to extract text from .txt files in the ZIP
            text_content.append("\nText file contents:")
            for filename in file_list[:10]:  # Only process first 10 files
                if filename.lower().endswith(('.txt', '.md', '.csv')):
                    try:
                        with zip_file.open(filename) as file:
                            content = file.read().decode('utf-8', errors='ignore')
                            text_content.append(f"\n--- {filename} ---")
                            text_content.append(content[:1000])  # First 1000 chars only
                    except:
                        continue
        
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"Failed to read ZIP file: {str(e)}")


def find_table_start(sheet) -> tuple:
    """Find where structured table data starts by looking for header patterns"""
    
    # Look for rows that appear to be table headers
    for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
        if not row or not any(row):
            continue
            
        row_values = [str(cell).strip() if cell is not None else "" for cell in row]
        non_empty_values = [val for val in row_values if val]
        
        # Skip if too few columns
        if len(non_empty_values) < 2:
            continue
        
        # Check if this looks like a header row
        if is_likely_header_row(non_empty_values):
            # Check if next few rows have data (confirming this is a table)
            if has_data_rows_following(sheet, row_idx, len(non_empty_values)):
                return row_idx, non_empty_values
    
    return None, None

def is_likely_header_row(row_values: list) -> bool:
    """Determine if a row looks like table headers"""
    
    # Common header patterns
    common_headers = [
        'name', 'id', 'number', 'phone', 'mobile', 'email', 'address', 
        'city', 'state', 'pincode', 'salary', 'age', 'date', 'amount',
        'description', 'type', 'status', 'code', 'value'
    ]
    
    # Check if row contains typical header words
    header_matches = 0
    for value in row_values:
        value_lower = value.lower().strip()
        if any(header in value_lower for header in common_headers):
            header_matches += 1
    
    # If more than half the columns look like headers
    if header_matches >= len(row_values) * 0.5:
        return True
    
    # Check for consistent formatting (short, capitalized words)
    short_caps_count = 0
    for value in row_values:
        if len(value) < 20 and (value.isupper() or value.istitle()):
            short_caps_count += 1
    
    return short_caps_count >= len(row_values) * 0.6

def has_data_rows_following(sheet, header_row_idx: int, expected_columns: int) -> bool:
    """Check if there are data rows following the potential header"""
    
    data_rows_found = 0
    max_check_rows = 10  # Check next 10 rows
    
    for row_idx in range(header_row_idx + 1, min(header_row_idx + max_check_rows + 1, sheet.max_row + 1)):
        row_values = []
        for cell in sheet[row_idx]:
            if cell.value is not None:
                row_values.append(str(cell.value).strip())
            else:
                row_values.append("")
        
        # Count non-empty cells
        non_empty = [val for val in row_values[:expected_columns] if val]
        
        # If row has data in at least half the expected columns
        if len(non_empty) >= expected_columns * 0.4:
            data_rows_found += 1
            
            # If we find 2+ data rows, it's likely a table
            if data_rows_found >= 2:
                return True
    
    return False

def is_empty_or_header_row(row_data: list, headers: list) -> bool:
    """Check if row is empty or duplicate header"""
    
    # Empty row
    if not any(val.strip() for val in row_data):
        return True
    
    # Check if it's a duplicate header row
    matches = 0
    for val, header in zip(row_data, headers):
        if val.strip().lower() == header.lower():
            matches += 1
    
    return matches >= len(headers) * 0.7  # 70% match with headers   

def chunk_excel_data(excel_data: dict, chunk_size: int = 5) -> list:
    """Chunk Excel structured data row-wise with column names per sheet"""
    chunks = []
    
    # First, add supporting text as separate chunks if exists
    if excel_data.get("supporting_text"):
        # supporting_text is a list, so join it first
        supporting_text_string = "\n".join(excel_data["supporting_text"])
        supporting_chunks = chunk_text_enhanced(supporting_text_string, chunk_size=200, overlap=20)
        for chunk in supporting_chunks:
            chunks.append(f"Supporting Information:\n{chunk}")
    
    # Then add structured data chunks
    current_sheet = None
    current_chunk_rows = []
    current_headers = []
    
    for row in excel_data["structured_rows"]:
        if current_sheet != row["sheet"]:
            # Save previous chunk if exists
            if current_chunk_rows:
                headers_line = "Columns: " + " | ".join(current_headers)
                chunk_text = f"Structured Data - Sheet: {current_sheet}\n{headers_line}\n" + "\n".join(current_chunk_rows)
                chunks.append(chunk_text)
            
            # Start new sheet
            current_sheet = row["sheet"]
            current_headers = row.get("headers", [])
            
            # Extract just the values from row_text
            row_text = str(row["row_text"])  # Ensure it's a string
            values_only = []
            for part in row_text.split(" | "):
                if ": " in part:
                    value = part.split(": ", 1)[1]
                    values_only.append(value)
            
            current_chunk_rows = [" | ".join(values_only)]
        else:
            # Extract just the values from row_text
            row_text = str(row["row_text"])  # Ensure it's a string
            values_only = []
            for part in row_text.split(" | "):
                if ": " in part:
                    value = part.split(": ", 1)[1]
                    values_only.append(value)
            
            current_chunk_rows.append(" | ".join(values_only))
            
            # Create chunk when reaching chunk_size
            if len(current_chunk_rows) >= chunk_size:
                headers_line = "Columns: " + " | ".join(current_headers)
                chunk_text = f"Structured Data - Sheet: {current_sheet}\n{headers_line}\n" + "\n".join(current_chunk_rows)
                chunks.append(chunk_text)
                current_chunk_rows = []
    
    # Add final chunk
    if current_chunk_rows and current_headers:
        headers_line = "Columns: " + " | ".join(current_headers)
        chunk_text = f"Structured Data - Sheet: {current_sheet}\n{headers_line}\n" + "\n".join(current_chunk_rows)
        chunks.append(chunk_text)
    
    return chunks

def chunk_text_smart(text: str, extraction_result: dict = None, chunk_size: int = 300, overlap: int = 50) -> list:
    """Smart chunking that handles Excel differently"""
    
    # Check if it's Excel data
    if extraction_result and extraction_result.get("file_type") == "xlsx" and extraction_result.get("excel_data"):
        print("📊 Using row-wise chunking for Excel file")
        return chunk_excel_data(extraction_result["excel_data"], chunk_size=5)  # 5 rows per chunk
    else:
        # Use regular word-based chunking for other files
        return chunk_text_enhanced(text, chunk_size, overlap)
    
def chunk_text(text: str, chunk_size: int = 200, overlap: int = 50) -> list:
    """Original chunking function for backward compatibility"""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks